from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.autoshape import TextFrame, AutoShapeType
from pptx.shapes.placeholder import SlidePlaceholder

# فایل template.pptx را بارگذاری می‌کنیم
# template_path = 'template.pptx'
# prs = Presentation(template_path)
#
# # برای هر اسلاید در فایل template
# for slide in prs.slides:
#     for shape in slide.shapes:
#         if isinstance(shape, TextFrame):
#             print("این یک TextBox است.")
#         elif isinstance(shape, Picture):
#             print("این یک Picture است.")
#         elif isinstance(shape, SlidePlaceholder):
#             print("این یک SlidePlaceholder است.")
#         elif isinstance(shape, AutoShapeType):
#             print("این یک AutoShape است.")
#         else:
#             print(f"نوع ناشناخته: {shape.__class__}")
#
# # فایل جدید با نام جدید ذخیره می‌کنیم
# output_path = 'output.pptx'
# prs.save(output_path)

questionnaire = {'title':'مقایسه دو پلتفرم اسنپ و تپسی',
                 'questions':[{'text': 'کدام بهتر است؟'},{'text': 'اسنپ را چقدر دوست دارید؟'}]}

class PptxElementDetector:
    elements = ['Picture','TextFrame','AutoShapeType','SlidePlaceholder', 'Shape']
    template_text = {
        'first_page': 'گزارش نظرسنجی {} در سنجاپ',
        'question_header': 'نظرسنجی انجام شده {} در سنجاپ'
    }
    def __init__(self, template_path):
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.new_prs = Presentation()

    def get_first_slide(self):
        return self.prs.slides[0]

    def get_questions_slide(self):
        return self.prs.slides[1]

    def get_final_slide(self):
        return self.prs.slides[-1]

    def set_first_slide_data(self,slide, title):
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                self.set_text(shape,title)

    def first_slide_layout(self):
        return self.prs.slides[0].slide_layout
    
    def question_slide_layout(self):
        return self.prs.slides[1].slide_layout

    def add_first_slide(self,title):
        slide = self.new_prs.slides.add_slide(self.first_slide_layout())
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                self.set_text(shape,title)
        output_path = 'output.pptx'
        self.prs.save(output_path)
    
    def add_question_slide(self,question):
        slide = self.new_prs.slides.add_slide(self.question_slide_layout())
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                self.set_text(shape,question)
        output_path = 'output.pptx'
        self.prs.save(output_path)

    def set_header_question_slide(self,slide ,header):
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                self.set_text(shape,header)

    def set_text(self, shape, text):
        if hasattr(shape, 'text'):
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    font_name = font.name
                    font_size = font.size
                    font_bold = font.bold
                    font_italic = font.italic

                    run.text = text

                    run.font.name = font_name
                    run.font.size = font_size
                    run.font.bold = font_bold
                    run.font.italic = font_italic

        # output_path = 'output.pptx'
        # self.prs.save(output_path)


    def create_slides(self):
        
        self.set_first_slide_data(self.get_first_slide(),self.template_text['first_page'].replace('{}',questionnaire['title']))
        for question in questionnaire['questions']:
            slide = self.get_questions_slide()
            slide_layout = slide.slide_layout
            new_slide = self.prs.slides.add_slide(slide_layout)
            self.set_header_question_slide(new_slide,question['text'])
        output_path = 'output.pptx'
        self.prs.save(output_path)



    # def detect(self):
    #     for slide in self.prs.slides:
    #         for shape in slide.shapes:
    #             print(shape.__class__)
    #             if shape.__class__.__name__ in self.elements:
    #                 if hasattr(shape,'text'):
    #                     self.set_text(shape, 'دکتر خان')
    #                 if shape.has_chart:
    #                     chart = shape.chart  # دسترسی به نمودار
    #                     print("نمودار پیدا شد!")
    #                     print(f"نوع نمودار: {chart.chart_type}")
    #                     # اینجا می‌توانید ویژگی‌های دیگری از نمودار را بررسی کنید یا تغییر دهید
    #                     # مثلا، دسترسی به داده‌های نمودار:
    #                     chart_data = chart.plots
    #                     for series in chart_data:
    #                         print(f"نام سری داده: {series.name}")
    #                         # داده‌های سری‌های مختلف را می‌توانید چاپ یا ویرایش کنید
    #                         for point in series.points:
    #                             print(f"برچسب: {point.label}, مقدار: {point.value}")

# result = PptxElementDetector('template.pptx').create_slides()
result = PptxElementDetector('template.pptx').add_first_slide('fdsafljalsd')
