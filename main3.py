from pptx.util import Inches
from pptx import Presentation
import shutil
import os
input_path = "template.pptx"

# Reload the presentation cleanly
prs = Presentation(input_path)
new_prs = Presentation('template.pptx')
questionnaire = {
    'title': 'مقایسه دو پلتفرم اسنپ و تپسی',
    'questions': [
        {'text': 'کدام بهتر است؟'},
        {'text': 'اسنپ را چقدر دوست دارید؟'}
    ]
}

# Copy the first slide (keeping layout and content)
first_slide_layout = prs.slides[0].slide_layout
first_slide = new_prs.slides.add_slide(first_slide_layout)
for shape in prs.slides[0].shapes:
    if shape.has_text_frame:
        new_shape = first_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
        new_shape.text_frame.text = shape.text

# Get layout of middle slide
middle_slide = prs.slides[1]
middle_layout = middle_slide.slide_layout

# Generate question slides
for q in questionnaire['questions']:
    slide = new_prs.slides.add_slide(middle_layout)
    for shape in middle_slide.shapes:
        if shape.has_text_frame:
            new_shape = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            if "?" in shape.text or "؟" in shape.text:
                new_shape.text_frame.text = q['text']
            else:
                new_shape.text_frame.text = shape.text

# Copy the last slide
last_slide_layout = prs.slides[-1].slide_layout
last_slide = new_prs.slides.add_slide(last_slide_layout)
for shape in prs.slides[-1].shapes:
    if shape.has_text_frame:
        new_shape = last_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
        new_shape.text_frame.text = shape.text

# Save the new presentation
output_path = "generated_questionnaire.pptx"
new_prs.save(output_path)

