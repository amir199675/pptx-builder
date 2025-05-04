import copy

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


class ChartSlideBuilder:
    def __init__(self, template_path: str, layout_index: int = 0, question_slide: int = 1,
                 intro_slide: int = 0, outro_slide: int = -1):
        self.presentation = Presentation(template_path)
        self.layout = self.presentation.slide_layouts[layout_index]
        self.bg_slide = self.presentation.slides[question_slide]
        self.bg_intro_slide = self.presentation.slides[intro_slide]
        self.bg_outro_slide = self.presentation.slides[outro_slide]

    def add_column_chart_slide(self, title: str, categories: list, values: list, series_name: str = "سری ۱"):
        # ساخت اسلاید
        slide = self.presentation.slides.add_slide(self.bg_slide.slide_layout)
        self._copy_background(slide)

        # افزودن عنوان
        if slide.shapes.title:
            slide.shapes.title.text = title

        # آماده‌سازی داده‌های چارت
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(series_name, values)

        # موقعیت چارت
        # x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4)
        x, y = Inches(0.5), Inches(1.5)
        cx, cy = Inches(9), Inches(5.5)
        # افزودن چارت
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

    def _copy_background(self, new_slide):
        for shape in self.bg_slide.shapes:
            if shape.shape_type != 1:  # exclude title box
                el = shape.element
                new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

    def add_pie_chart_slide(self, title, categories, values, series_name="سری ۱"):
        new_slide = self.presentation.slides.add_slide(self.bg_slide.slide_layout)
        self._copy_background(new_slide)

        if new_slide.shapes.title:
            new_slide.shapes.title.text = title

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(series_name, values)

        x, y = Inches(1.5), Inches(1.5)
        cx, cy = Inches(7), Inches(5.5)

        chart = new_slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        # ✅ اینجا فرق اصلیه:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_value = True
        data_labels.show_category_name = True
        # data_labels.show_percentage = True

    def add_intro_slide(self, text):
        # layout = self.presentation.slide_layouts[layout_index]  # مثلاً Title layout
        slide = self.presentation.slides.add_slide(self.bg_intro_slide.slide_layout)
        self._copy_background(slide)
        self.add_custom_text_slide(slide, text, 3, 0.5)

    def add_outro_slide(self, text):
        # layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(self.bg_outro_slide.slide_layout)
        self._copy_background(slide)
        self.add_custom_text_slide(slide, text, 4, 5)

    def copy_shapes(self, from_slide, to_slide):
        for shape in from_slide.shapes:
            el = shape.element
            new_el = copy.deepcopy(el)
            to_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    def add_custom_text_slide(self, slide, text, x=1, y=2, width=8, height=2, font_size=32,font_name='Calibri'):
        # slide = self.presentation.slides.add_slide(self.bg_slide.slide_layout)
        # self._copy_background(slide)

        textbox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = font_name

    def remove_slide(self, index: int):
        slides = self.presentation.slides._sldIdLst  # لیست XML اسلایدها
        slide_id = list(slides)[index]
        slides.remove(slide_id)

    def save(self, output_path: str):
        self.presentation.save(output_path)

    def create_pptx_file_from_data(self,data):
        self.add_intro_slide("به نام خدا")

        for i in data:
            chart_categories = []
            chart_values = []
            for option in i['options']:
                chart_categories.append(option['value_title'])
                chart_values.append(option['percentage'])
            self.add_pie_chart_slide(
                title=i['question_title'],
                categories=chart_categories,
                values=chart_values,
                series_name=i['question_title']
            )
            self.add_column_chart_slide(
                title=i['question_title'],
                categories=chart_categories,
                values=chart_values,
                series_name=i['question_title']
            )
        self.add_outro_slide("پایان ارائه")
        for i in range(4):
            self.remove_slide(0)
        self.save("output_real_chart_oop.pptx")




data = [
    {
        "question_order": 2,
        "question_title": "سطح تحصیلات خود را مشخص بفرمایید",
        "options": [
            {
                "value_id": 29612,
                "value_title": "زیر دیپلم",
                "count": 41,
                "percentage": "2.12",
                "women_count": 18,
                "men_count": 23,
                "non_available": 0
            },
            {
                "value_id": 29613,
                "value_title": "دیپلم",
                "count": 93,
                "percentage": "4.80",
                "women_count": 38,
                "men_count": 55,
                "non_available": 0
            },
            {
                "value_id": 29614,
                "value_title": "کاردانی",
                "count": 21,
                "percentage": "1.08",
                "women_count": 8,
                "men_count": 13,
                "non_available": 0
            },
            {
                "value_id": 29615,
                "value_title": "کارشناسی",
                "count": 44,
                "percentage": "2.27",
                "women_count": 18,
                "men_count": 26,
                "non_available": 0
            },
            {
                "value_id": 29616,
                "value_title": "کارشناسی ارشد",
                "count": 15,
                "percentage": "0.77",
                "women_count": 2,
                "men_count": 13,
                "non_available": 0
            },
            {
                "value_id": 29617,
                "value_title": "دکتری و بالاتر",
                "count": 5,
                "percentage": "0.26",
                "women_count": 2,
                "men_count": 3,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 3,
        "question_title": "دسته شغلی خود را مشخص بفرمایید",
        "options": [
            {
                "value_id": 29626,
                "value_title": "شاغل بخش خصوصی",
                "count": 23,
                "percentage": "1.19",
                "women_count": 3,
                "men_count": 20,
                "non_available": 0
            },
            {
                "value_id": 29627,
                "value_title": "شاغل بخش دولتی",
                "count": 17,
                "percentage": "0.88",
                "women_count": 2,
                "men_count": 15,
                "non_available": 0
            },
            {
                "value_id": 29628,
                "value_title": "بازنشسته",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29629,
                "value_title": "آزاد",
                "count": 50,
                "percentage": "2.58",
                "women_count": 5,
                "men_count": 45,
                "non_available": 0
            },
            {
                "value_id": 29630,
                "value_title": "دانشجو",
                "count": 35,
                "percentage": "1.81",
                "women_count": 15,
                "men_count": 20,
                "non_available": 0
            },
            {
                "value_id": 29631,
                "value_title": "خانه‌دار",
                "count": 44,
                "percentage": "2.27",
                "women_count": 41,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29632,
                "value_title": "سرباز",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29633,
                "value_title": "بیکار",
                "count": 30,
                "percentage": "1.55",
                "women_count": 14,
                "men_count": 16,
                "non_available": 0
            },
            {
                "value_id": 29634,
                "value_title": "سایر",
                "count": 12,
                "percentage": "0.62",
                "women_count": 4,
                "men_count": 8,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 5,
        "question_title": "آیا شما بیسکوئیت وین جوی کشمشی برند زر کام در شش ماه گذشته استفاده کرده‌اید؟",
        "options": [
            {
                "value_id": 29635,
                "value_title": "خیر",
                "count": 138,
                "percentage": "7.13",
                "women_count": 56,
                "men_count": 82,
                "non_available": 0
            },
            {
                "value_id": 29636,
                "value_title": "بله",
                "count": 78,
                "percentage": "4.03",
                "women_count": 30,
                "men_count": 48,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 6,
        "question_title": "میزان رضایت شما از طعم و مزه این محصول چقدر است؟",
        "options": [
            {
                "value_id": 29637,
                "value_title": "1",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29639,
                "value_title": "3",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29640,
                "value_title": "4",
                "count": 1,
                "percentage": "0.05",
                "women_count": 1,
                "men_count": 0,
                "non_available": 0
            },
            {
                "value_id": 29641,
                "value_title": "5",
                "count": 3,
                "percentage": "0.15",
                "women_count": 0,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29642,
                "value_title": "6",
                "count": 15,
                "percentage": "0.77",
                "women_count": 5,
                "men_count": 10,
                "non_available": 0
            },
            {
                "value_id": 29643,
                "value_title": "7",
                "count": 9,
                "percentage": "0.46",
                "women_count": 4,
                "men_count": 5,
                "non_available": 0
            },
            {
                "value_id": 29644,
                "value_title": "8",
                "count": 12,
                "percentage": "0.62",
                "women_count": 2,
                "men_count": 10,
                "non_available": 0
            },
            {
                "value_id": 29645,
                "value_title": "9",
                "count": 6,
                "percentage": "0.31",
                "women_count": 2,
                "men_count": 4,
                "non_available": 0
            },
            {
                "value_id": 29646,
                "value_title": "10",
                "count": 31,
                "percentage": "1.60",
                "women_count": 15,
                "men_count": 16,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 7,
        "question_title": "میزان رضایت شما از طعم و مزه کشمش این محصول چقدر است؟",
        "options": [
            {
                "value_id": 29720,
                "value_title": "4",
                "count": 5,
                "percentage": "0.26",
                "women_count": 2,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29721,
                "value_title": "5",
                "count": 10,
                "percentage": "0.52",
                "women_count": 4,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29722,
                "value_title": "6",
                "count": 12,
                "percentage": "0.62",
                "women_count": 3,
                "men_count": 9,
                "non_available": 0
            },
            {
                "value_id": 29723,
                "value_title": "7",
                "count": 8,
                "percentage": "0.41",
                "women_count": 2,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29724,
                "value_title": "8",
                "count": 9,
                "percentage": "0.46",
                "women_count": 3,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29725,
                "value_title": "9",
                "count": 13,
                "percentage": "0.67",
                "women_count": 6,
                "men_count": 7,
                "non_available": 0
            },
            {
                "value_id": 29726,
                "value_title": "10",
                "count": 22,
                "percentage": "1.14",
                "women_count": 10,
                "men_count": 12,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 8,
        "question_title": "میزان رضایت شما از وضعیت ظاهری (له‏ شدگی، ترک‏ خوردگی، شکستگی، دوییدگی) این محصول چقدر است؟",
        "options": [
            {
                "value_id": 29647,
                "value_title": "1",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29648,
                "value_title": "2",
                "count": 1,
                "percentage": "0.05",
                "women_count": 0,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29649,
                "value_title": "3",
                "count": 1,
                "percentage": "0.05",
                "women_count": 0,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29650,
                "value_title": "4",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29651,
                "value_title": "5",
                "count": 6,
                "percentage": "0.31",
                "women_count": 3,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29652,
                "value_title": "6",
                "count": 17,
                "percentage": "0.88",
                "women_count": 4,
                "men_count": 13,
                "non_available": 0
            },
            {
                "value_id": 29653,
                "value_title": "7",
                "count": 9,
                "percentage": "0.46",
                "women_count": 3,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29654,
                "value_title": "8",
                "count": 10,
                "percentage": "0.52",
                "women_count": 4,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29655,
                "value_title": "9",
                "count": 10,
                "percentage": "0.52",
                "women_count": 4,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29656,
                "value_title": "10",
                "count": 19,
                "percentage": "0.98",
                "women_count": 10,
                "men_count": 9,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 9,
        "question_title": "میزان رضایت شما از سهولت مصرف بسته ‏بندی این محصول چقدر است؟",
        "options": [
            {
                "value_id": 29708,
                "value_title": "2",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29709,
                "value_title": "3",
                "count": 1,
                "percentage": "0.05",
                "women_count": 0,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29710,
                "value_title": "4",
                "count": 3,
                "percentage": "0.15",
                "women_count": 2,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29711,
                "value_title": "5",
                "count": 1,
                "percentage": "0.05",
                "women_count": 0,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29712,
                "value_title": "6",
                "count": 15,
                "percentage": "0.77",
                "women_count": 4,
                "men_count": 11,
                "non_available": 0
            },
            {
                "value_id": 29713,
                "value_title": "7",
                "count": 12,
                "percentage": "0.62",
                "women_count": 2,
                "men_count": 10,
                "non_available": 0
            },
            {
                "value_id": 29714,
                "value_title": "8",
                "count": 9,
                "percentage": "0.46",
                "women_count": 2,
                "men_count": 7,
                "non_available": 0
            },
            {
                "value_id": 29715,
                "value_title": "9",
                "count": 8,
                "percentage": "0.41",
                "women_count": 5,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29716,
                "value_title": "10",
                "count": 27,
                "percentage": "1.39",
                "women_count": 13,
                "men_count": 14,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 10,
        "question_title": "میزان رضایت شما از کیفیت دوخت بسته‏ بندی چقدر است؟",
        "options": [
            {
                "value_id": 29748,
                "value_title": "2",
                "count": 1,
                "percentage": "0.05",
                "women_count": 0,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29750,
                "value_title": "4",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29751,
                "value_title": "5",
                "count": 5,
                "percentage": "0.26",
                "women_count": 2,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29752,
                "value_title": "6",
                "count": 13,
                "percentage": "0.67",
                "women_count": 5,
                "men_count": 8,
                "non_available": 0
            },
            {
                "value_id": 29753,
                "value_title": "7",
                "count": 12,
                "percentage": "0.62",
                "women_count": 4,
                "men_count": 8,
                "non_available": 0
            },
            {
                "value_id": 29754,
                "value_title": "8",
                "count": 12,
                "percentage": "0.62",
                "women_count": 2,
                "men_count": 10,
                "non_available": 0
            },
            {
                "value_id": 29755,
                "value_title": "9",
                "count": 12,
                "percentage": "0.62",
                "women_count": 6,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29756,
                "value_title": "10",
                "count": 21,
                "percentage": "1.08",
                "women_count": 9,
                "men_count": 12,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 11,
        "question_title": "میزان رضایت شما از نظر اطلاعات درج ‏شده روی محصول بر روی بسته ‏بندی چقدر است؟",
        "options": [
            {
                "value_id": 29660,
                "value_title": "4",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29661,
                "value_title": "5",
                "count": 5,
                "percentage": "0.26",
                "women_count": 1,
                "men_count": 4,
                "non_available": 0
            },
            {
                "value_id": 29662,
                "value_title": "6",
                "count": 12,
                "percentage": "0.62",
                "women_count": 5,
                "men_count": 7,
                "non_available": 0
            },
            {
                "value_id": 29663,
                "value_title": "7",
                "count": 13,
                "percentage": "0.67",
                "women_count": 3,
                "men_count": 10,
                "non_available": 0
            },
            {
                "value_id": 29664,
                "value_title": "8",
                "count": 10,
                "percentage": "0.52",
                "women_count": 4,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29665,
                "value_title": "9",
                "count": 12,
                "percentage": "0.62",
                "women_count": 5,
                "men_count": 7,
                "non_available": 0
            },
            {
                "value_id": 29666,
                "value_title": "10",
                "count": 26,
                "percentage": "1.34",
                "women_count": 12,
                "men_count": 14,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 12,
        "question_title": "میزان رضایت شما از  طرح بسته بندی چقدر است؟",
        "options": [
            {
                "value_id": 29739,
                "value_title": "3",
                "count": 2,
                "percentage": "0.10",
                "women_count": 0,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29740,
                "value_title": "4",
                "count": 1,
                "percentage": "0.05",
                "women_count": 0,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29741,
                "value_title": "5",
                "count": 6,
                "percentage": "0.31",
                "women_count": 0,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29742,
                "value_title": "6",
                "count": 14,
                "percentage": "0.72",
                "women_count": 5,
                "men_count": 9,
                "non_available": 0
            },
            {
                "value_id": 29743,
                "value_title": "7",
                "count": 10,
                "percentage": "0.52",
                "women_count": 5,
                "men_count": 5,
                "non_available": 0
            },
            {
                "value_id": 29744,
                "value_title": "8",
                "count": 14,
                "percentage": "0.72",
                "women_count": 5,
                "men_count": 9,
                "non_available": 0
            },
            {
                "value_id": 29745,
                "value_title": "9",
                "count": 6,
                "percentage": "0.31",
                "women_count": 3,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29746,
                "value_title": "10",
                "count": 27,
                "percentage": "1.39",
                "women_count": 13,
                "men_count": 14,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 13,
        "question_title": "میزان رضایت شما از دسترسی به تنوع این محصولات از نظر طعم (وانیل و ...) چقدر است؟",
        "options": [
            {
                "value_id": 29729,
                "value_title": "3",
                "count": 1,
                "percentage": "0.05",
                "women_count": 1,
                "men_count": 0,
                "non_available": 0
            },
            {
                "value_id": 29730,
                "value_title": "4",
                "count": 5,
                "percentage": "0.26",
                "women_count": 1,
                "men_count": 4,
                "non_available": 0
            },
            {
                "value_id": 29731,
                "value_title": "5",
                "count": 6,
                "percentage": "0.31",
                "women_count": 2,
                "men_count": 4,
                "non_available": 0
            },
            {
                "value_id": 29732,
                "value_title": "6",
                "count": 13,
                "percentage": "0.67",
                "women_count": 4,
                "men_count": 9,
                "non_available": 0
            },
            {
                "value_id": 29733,
                "value_title": "7",
                "count": 11,
                "percentage": "0.57",
                "women_count": 1,
                "men_count": 10,
                "non_available": 0
            },
            {
                "value_id": 29734,
                "value_title": "8",
                "count": 9,
                "percentage": "0.46",
                "women_count": 3,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29735,
                "value_title": "9",
                "count": 8,
                "percentage": "0.41",
                "women_count": 5,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29736,
                "value_title": "10",
                "count": 26,
                "percentage": "1.34",
                "women_count": 13,
                "men_count": 13,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 14,
        "question_title": "میزان رضایت شما از دسترسی به انواع محصولات (وین‏جوی و ....) چقدر است؟",
        "options": [
            {
                "value_id": 29698,
                "value_title": "2",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29699,
                "value_title": "3",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29700,
                "value_title": "4",
                "count": 4,
                "percentage": "0.21",
                "women_count": 1,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29701,
                "value_title": "5",
                "count": 5,
                "percentage": "0.26",
                "women_count": 2,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29702,
                "value_title": "6",
                "count": 9,
                "percentage": "0.46",
                "women_count": 4,
                "men_count": 5,
                "non_available": 0
            },
            {
                "value_id": 29703,
                "value_title": "7",
                "count": 13,
                "percentage": "0.67",
                "women_count": 1,
                "men_count": 12,
                "non_available": 0
            },
            {
                "value_id": 29704,
                "value_title": "8",
                "count": 16,
                "percentage": "0.83",
                "women_count": 7,
                "men_count": 9,
                "non_available": 0
            },
            {
                "value_id": 29705,
                "value_title": "9",
                "count": 9,
                "percentage": "0.46",
                "women_count": 6,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29706,
                "value_title": "10",
                "count": 20,
                "percentage": "1.03",
                "women_count": 8,
                "men_count": 12,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 15,
        "question_title": "میزان رضایت شما از قیمت این محصول نسبت به محصولات سایر رقبا چقدر است؟",
        "options": [
            {
                "value_id": 29667,
                "value_title": "1",
                "count": 2,
                "percentage": "0.10",
                "women_count": 0,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29669,
                "value_title": "3",
                "count": 2,
                "percentage": "0.10",
                "women_count": 0,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29670,
                "value_title": "4",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29671,
                "value_title": "5",
                "count": 6,
                "percentage": "0.31",
                "women_count": 3,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29672,
                "value_title": "6",
                "count": 20,
                "percentage": "1.03",
                "women_count": 7,
                "men_count": 13,
                "non_available": 0
            },
            {
                "value_id": 29673,
                "value_title": "7",
                "count": 9,
                "percentage": "0.46",
                "women_count": 2,
                "men_count": 7,
                "non_available": 0
            },
            {
                "value_id": 29674,
                "value_title": "8",
                "count": 9,
                "percentage": "0.46",
                "women_count": 4,
                "men_count": 5,
                "non_available": 0
            },
            {
                "value_id": 29675,
                "value_title": "9",
                "count": 6,
                "percentage": "0.31",
                "women_count": 3,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29676,
                "value_title": "10",
                "count": 24,
                "percentage": "1.24",
                "women_count": 11,
                "men_count": 13,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 16,
        "question_title": "میزان تمایل شما به خرید مجدد این محصول چقدر است؟",
        "options": [
            {
                "value_id": 29677,
                "value_title": "1",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29680,
                "value_title": "4",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29681,
                "value_title": "5",
                "count": 3,
                "percentage": "0.15",
                "women_count": 0,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29682,
                "value_title": "6",
                "count": 18,
                "percentage": "0.93",
                "women_count": 7,
                "men_count": 11,
                "non_available": 0
            },
            {
                "value_id": 29683,
                "value_title": "7",
                "count": 15,
                "percentage": "0.77",
                "women_count": 4,
                "men_count": 11,
                "non_available": 0
            },
            {
                "value_id": 29684,
                "value_title": "8",
                "count": 13,
                "percentage": "0.67",
                "women_count": 7,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29685,
                "value_title": "9",
                "count": 5,
                "percentage": "0.26",
                "women_count": 2,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29686,
                "value_title": "10",
                "count": 22,
                "percentage": "1.14",
                "women_count": 9,
                "men_count": 13,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 17,
        "question_title": "میزان تمایل شما به تشویق دیگران به خرید این محصول چقدر است؟",
        "options": [
            {
                "value_id": 29687,
                "value_title": "1",
                "count": 2,
                "percentage": "0.10",
                "women_count": 1,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29688,
                "value_title": "2",
                "count": 1,
                "percentage": "0.05",
                "women_count": 1,
                "men_count": 0,
                "non_available": 0
            },
            {
                "value_id": 29689,
                "value_title": "3",
                "count": 1,
                "percentage": "0.05",
                "women_count": 0,
                "men_count": 1,
                "non_available": 0
            },
            {
                "value_id": 29691,
                "value_title": "5",
                "count": 7,
                "percentage": "0.36",
                "women_count": 1,
                "men_count": 6,
                "non_available": 0
            },
            {
                "value_id": 29692,
                "value_title": "6",
                "count": 13,
                "percentage": "0.67",
                "women_count": 4,
                "men_count": 9,
                "non_available": 0
            },
            {
                "value_id": 29693,
                "value_title": "7",
                "count": 14,
                "percentage": "0.72",
                "women_count": 3,
                "men_count": 11,
                "non_available": 0
            },
            {
                "value_id": 29694,
                "value_title": "8",
                "count": 14,
                "percentage": "0.72",
                "women_count": 6,
                "men_count": 8,
                "non_available": 0
            },
            {
                "value_id": 29695,
                "value_title": "9",
                "count": 5,
                "percentage": "0.26",
                "women_count": 3,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29696,
                "value_title": "10",
                "count": 24,
                "percentage": "1.24",
                "women_count": 12,
                "men_count": 12,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 18,
        "question_title": "آیا تا کنون در هنگام باز کردن این محصول، حشرات موذی در محصول مشاهده کرده اید؟",
        "options": [
            {
                "value_id": 29620,
                "value_title": "بله",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29621,
                "value_title": "خیر",
                "count": 78,
                "percentage": "4.03",
                "women_count": 30,
                "men_count": 48,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 19,
        "question_title": "آیا تا کنون در هنگام باز کردن این محصول، جسم خارجی مانند: ذرات فلزی، مو، پلاستیک، نخ گونی، چسب نواری و ... در محصول مشاهده کرده اید؟",
        "options": [
            {
                "value_id": 29622,
                "value_title": "بله",
                "count": 3,
                "percentage": "0.15",
                "women_count": 1,
                "men_count": 2,
                "non_available": 0
            },
            {
                "value_id": 29623,
                "value_title": "خیر",
                "count": 79,
                "percentage": "4.08",
                "women_count": 30,
                "men_count": 49,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 20,
        "question_title": "آیا تا کنون بر اثر مصرف این محصول، دچار مسمومیت غذایی شده اید؟",
        "options": [
            {
                "value_id": 29625,
                "value_title": "خیر",
                "count": 82,
                "percentage": "4.24",
                "women_count": 31,
                "men_count": 51,
                "non_available": 0
            }
        ]
    },
    {
        "question_order": 21,
        "question_title": "آیا تا کنون در هنگام باز کردن این محصول، کپک‏زدگی مشاهده کرده اید؟",
        "options": [
            {
                "value_id": 29618,
                "value_title": "بله",
                "count": 6,
                "percentage": "0.31",
                "women_count": 3,
                "men_count": 3,
                "non_available": 0
            },
            {
                "value_id": 29619,
                "value_title": "خیر",
                "count": 76,
                "percentage": "3.93",
                "women_count": 28,
                "men_count": 48,
                "non_available": 0
            }
        ]
    }
]

builder = ChartSlideBuilder("template.pptx")
builder.create_pptx_file_from_data(data)