import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os

# بارگذاری قالب پاورپوینت
prs = Presentation("template.pptx")

# انتخاب layout مورد نظر از قالب
custom_layout = prs.slide_layouts[0]  # مثلا "Title and Content"

# داده‌ی شبیه‌سازی‌شده JSON
data = {
    "فروش ماهانه": {
        "فروردین": 120,
        "اردیبهشت": 150,
        "خرداد": 180
    },
    "هزینه ماهانه": {
        "فروردین": 100,
        "اردیبهشت": 110,
        "خرداد": 130
    }
}

for title, values in data.items():
    # ساخت نمودار با matplotlib
    fig, ax = plt.subplots()
    ax.bar(values.keys(), values.values())
    ax.set_title(title)
    chart_path = f"{title}.png"
    plt.savefig(chart_path, bbox_inches='tight')
    plt.close()

    # ساخت اسلاید با قالب اختصاصی
    slide = prs.slides.add_slide(custom_layout)

    # درج عنوان
    if slide.shapes.title:
        slide.shapes.title.text = title

    # درج تصویر نمودار
    slide.shapes.add_picture(chart_path, Inches(1), Inches(2), width=Inches(6))

    # حذف فایل موقت تصویر
    os.remove(chart_path)

# ذخیره پاورپوینت نهایی
prs.save("output_with_template.pptx")