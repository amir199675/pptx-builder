# import json
# import matplotlib.pyplot as plt
# from pptx import Presentation
# from pptx.util import Inches
# import os
#
# # داده‌ی JSON
# data = {
#     "فروش ماهانه": {
#         "فروردین": 120,
#         "اردیبهشت": 150,
#         "خرداد": 180
#     },
#     "هزینه ماهانه": {
#         "فروردین": 100,
#         "اردیبهشت": 110,
#         "خرداد": 130
#     }
# }
#
# # ساخت پاورپوینت
# prs = Presentation()
#
# for title, values in data.items():
#     # ساخت نمودار با matplotlib
#     fig, ax = plt.subplots()
#     ax.bar(values.keys(), values.values())
#     ax.set_title(title)
#     chart_path = f"{title}.png"
#     plt.savefig(chart_path)
#     plt.close()
#
#     # افزودن اسلاید و درج چارت
#     slide = prs.slides.add_slide(prs.slide_layouts[5])
#     title_box = slide.shapes.title
#     title_box.text = title
#     slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(6))
#
#     # پاک‌سازی فایل تصویر
#     os.remove(chart_path)
#
# # ذخیره پاورپوینت
# prs.save("presentation_with_charts.pptx")



# import json
# import matplotlib.pyplot as plt
# from pptx import Presentation
# from pptx.util import Inches
# import os
#
# # بارگذاری قالب پاورپوینت
# prs = Presentation("template.pptx")
#
# # انتخاب layout مورد نظر از قالب
# custom_layout = prs.slide_layouts[0]  # مثلا "Title and Content"
#
# # داده‌ی شبیه‌سازی‌شده JSON
# data = {
#     "فروش ماهانه": {
#         "فروردین": 120,
#         "اردیبهشت": 150,
#         "خرداد": 180
#     },
#     "هزینه ماهانه": {
#         "فروردین": 100,
#         "اردیبهشت": 110,
#         "خرداد": 130
#     }
# }
#
# for title, values in data.items():
#     # ساخت نمودار با matplotlib
#     fig, ax = plt.subplots()
#     ax.bar(values.keys(), values.values())
#     ax.set_title(title)
#     chart_path = f"{title}.png"
#     plt.savefig(chart_path, bbox_inches='tight')
#     plt.close()
#
#     # ساخت اسلاید با قالب اختصاصی
#     slide = prs.slides.add_slide(custom_layout)
#
#     # درج عنوان
#     if slide.shapes.title:
#         slide.shapes.title.text = title
#
#     # درج تصویر نمودار
#     slide.shapes.add_picture(chart_path, Inches(1), Inches(2), width=Inches(6))
#
#     # حذف فایل موقت تصویر
#     os.remove(chart_path)
#
# # ذخیره پاورپوینت نهایی
# prs.save("output_with_template.pptx")


# بازنویسی مجدد پس از ریست شدن محیط
# import json
# from pptx import Presentation
# from pptx.util import Pt
# from pptx.chart.data import CategoryChartData
# import arabic_reshaper
# from bidi.algorithm import get_display
#
# # تابع اصلاح متن فارسی
# def fix_farsi_text(text):
#     reshaped_text = arabic_reshaper.reshape(text)
#     return get_display(reshaped_text)
#
# # داده نمونه
# json_data = [
#     {
#         "template_slide": 1,
#         "text": "سوال اول",
#         "chart_data": {
#             "categories": ["الف", "ب", "ج"],
#             "series": {
#                 "درصد پاسخ صحیح": [80, 60, 90]
#             }
#         }
#     },
#     {
#         "template_slide": 2,
#         "text": "نتیجه‌گیری"
#     }
# ]
#
# source_pptx = "template.pptx"
# destination_pptx = "output_4.pptx"
#
# src = Presentation(source_pptx)
# dst = Presentation()
# dst.slides._sldIdLst.clear()
# dst.slide_width = src.slide_width
# dst.slide_height = src.slide_height
#
# for entry in json_data:
#     src_slide_index = entry["template_slide"]
#     custom_text = fix_farsi_text(entry["text"])
#     chart_info = entry.get("chart_data")
#
#     template_slide = src.slides[src_slide_index]
#     layout = template_slide.slide_layout
#     new_slide = dst.slides.add_slide(layout)
#
#     for shape in template_slide.shapes:
#         if shape.has_text_frame:
#             textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
#             tf = textbox.text_frame
#             tf.text = custom_text
#             for p in shape.text_frame.paragraphs[1:]:
#                 tf.add_paragraph().text = fix_farsi_text(p.text)
#
#         elif shape.has_chart and chart_info:
#             chart = shape.chart
#             chart_data = CategoryChartData()
#             chart_data.categories = [fix_farsi_text(cat) for cat in chart_info["categories"]]
#             for series_name, values in chart_info["series"].items():
#                 chart_data.add_series(fix_farsi_text(series_name), values)
#             chart.replace_data(chart_data)
#
#         elif shape.shape_type == 13:  # Picture
#             image = shape.image
#             with open("temp_img.png", "wb") as f:
#                 f.write(image.blob)
#             new_slide.shapes.add_picture("temp_img.png", shape.left, shape.top, shape.width, shape.height)
#
# dst.save(destination_pptx)
# print(f"✅ فایل خروجی ساخته شد: {destination_pptx}")


from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# داده‌ها
data = {
    "فروش ماهانه": {
        "فروردین": 120,
        "اردیبهشت": 150,
        "خرداد": 180
    },
    "هزینه ماهانه": {
        "فروردین": 100,
        "اردیبهشت": 110,
        "خرداد": 130
    }
}

# ساخت فایل جدید پاورپوینت یا باز کردن قالب
prs = Presentation("template.pptx")
slide_layout = prs.slide_layouts[0]  # layout خالی یا content

for title, values in data.items():
    # ساخت اسلاید جدید
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    # ساخت داده چارت
    chart_data = CategoryChartData()
    chart_data.categories = list(values.keys())
    chart_data.add_series(f'{title}', list(values.values()))

    # محل چارت
    x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4)

    # افزودن چارت به اسلاید
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

# ذخیره فایل
prs.save("output_real_chart.pptx")
