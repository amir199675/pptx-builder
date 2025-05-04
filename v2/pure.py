from pptx import Presentation

# باز کردن فایل پاورپوینت
slide_mapper = {
    'first': 'Slide_001',
    'question_chart': 'Slide_002',
    'question_table': 'Slide_003',
    'last': 'Slide_004',
}
def add_note():
    prs = Presentation("template_3.pptx")
    for idx, slide in enumerate(prs.slides):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"ID: Slide_{idx+1:03}"

    prs.save("output_3.pptx")


def find_slide(target_id):
    prs = Presentation("output_3.pptx")

    for index, slide in enumerate(prs.slides):
        notes = slide.notes_slide.notes_text_frame.text
        if target_id in notes:
            print(f"اسلاید پیدا شد: شماره {index + 1}")
            return slide
    else:
        print("اسلاید مورد نظر پیدا نشد.")
        return None

def copy_prs():
    src = Presentation("template_3.pptx")
    dst = Presentation("output_3.pptx")

    src_slide = src.slides[0]
    layout = dst.slide_layouts[0]
    new_slide = dst.slides.add_slide(layout)

    for shape in src_slide.shapes:
        if shape.shape_type == 13:  # Picture
            image = shape.image
            with open("temp_img.png", "wb") as f:
                f.write(image.blob)
            new_slide.shapes.add_picture("temp_img.png", shape.left, shape.top, shape.width, shape.height)
        elif shape.has_text_frame:
            textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            textbox.text_frame.text = shape.text

    dst.save("destination_with_slide.pptx")

import zipfile
import os
import shutil
import tempfile
import xml.etree.ElementTree as ET

def copy_slide(src_path, dst_path, slide_index):  # slide_index = 1 → slide2.xml
    temp_src = tempfile.mkdtemp()
    temp_dst = tempfile.mkdtemp()

    # 1. unzip both files
    with zipfile.ZipFile(src_path, 'r') as zip_ref:
        zip_ref.extractall(temp_src)
    with zipfile.ZipFile(dst_path, 'r') as zip_ref:
        zip_ref.extractall(temp_dst)

    # 2. find next available slide number in dst
    slide_folder = os.path.join(temp_dst, 'ppt', 'slides')
    next_slide_num = len([f for f in os.listdir(slide_folder) if f.startswith('slide')]) + 1
    new_slide_name = f'slide{next_slide_num}.xml'

    # 3. copy slide file from src
    src_slide = f'slide{slide_index + 1}.xml'
    shutil.copy(os.path.join(temp_src, 'ppt', 'slides', src_slide),
                os.path.join(temp_dst, 'ppt', 'slides', new_slide_name))

    # 4. copy rels of slide
    src_slide_rels = os.path.join(temp_src, 'ppt', 'slides', '_rels', src_slide + '.rels')
    dst_slide_rels = os.path.join(temp_dst, 'ppt', 'slides', '_rels', new_slide_name + '.rels')
    if os.path.exists(src_slide_rels):
        shutil.copy(src_slide_rels, dst_slide_rels)

    # 5. add entry to presentation.xml
    pres_xml = os.path.join(temp_dst, 'ppt', 'presentation.xml')
    tree = ET.parse(pres_xml)
    root = tree.getroot()
    sldIdLst = root.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')

    # find last id and increment
    ids = [int(sld.get('id')) for sld in sldIdLst.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')]
    max_id = max(ids)
    r_id = f'rId{max_id + 256}'

    # add new slide node
    new_sld = ET.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
    new_sld.set('id', str(max_id + 1))
    new_sld.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', r_id)
    sldIdLst.append(new_sld)
    tree.write(pres_xml)

    # 6. add slide relationship to presentation.xml.rels
    rels_path = os.path.join(temp_dst, 'ppt', '_rels', 'presentation.xml.rels')
    tree = ET.parse(rels_path)
    root = tree.getroot()

    new_rel = ET.Element('Relationship')
    new_rel.set('Id', r_id)
    new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
    new_rel.set('Target', f'slides/{new_slide_name}')
    root.append(new_rel)
    tree.write(rels_path)

    # 7. rezip the destination folder
    new_pptx = dst_path.replace('.pptx', '_with_slide.pptx')
    with zipfile.ZipFile(new_pptx, 'w', zipfile.ZIP_DEFLATED) as pptx:
        for folder_name, subfolders, filenames in os.walk(temp_dst):
            for filename in filenames:
                abs_path = os.path.join(folder_name, filename)
                rel_path = os.path.relpath(abs_path, temp_dst)
                pptx.write(abs_path, rel_path)

    # 8. clean up
    shutil.rmtree(temp_src)
    shutil.rmtree(temp_dst)

    print(f"✅ اسلاید دوم با موفقیت به {new_pptx} منتقل شد.")

# استفاده:
# copy_slide("template_3.pptx", "output_3.pptx", 0)  # اسلاید دوم (اندیس 1)

# import zipfile
# import os
# import shutil
# import tempfile
# import xml.etree.ElementTree as ET
#
#
# def copy_slide_with_layout(src_path, dst_path, slide_index=1):
#     ET.register_namespace('', "http://schemas.openxmlformats.org/presentationml/2006/main")
#     ET.register_namespace('r', "http://schemas.openxmlformats.org/officeDocument/2006/relationships")
#
#     temp_src = tempfile.mkdtemp()
#     temp_dst = tempfile.mkdtemp()
#
#     # Unzip both files
#     with zipfile.ZipFile(src_path, 'r') as zip_ref:
#         zip_ref.extractall(temp_src)
#     with zipfile.ZipFile(dst_path, 'r') as zip_ref:
#         zip_ref.extractall(temp_dst)
#
#     slide_id = slide_index + 1
#     src_slide_name = f'slide{slide_id}.xml'
#     new_slide_num = len(os.listdir(os.path.join(temp_dst, 'ppt', 'slides'))) + 1
#     new_slide_name = f'slide{new_slide_num}.xml'
#
#     # Copy slide XML
#     shutil.copy(
#         os.path.join(temp_src, 'ppt', 'slides', src_slide_name),
#         os.path.join(temp_dst, 'ppt', 'slides', new_slide_name)
#     )
#
#     # Copy slide relationships
#     src_slide_rels = os.path.join(temp_src, 'ppt', 'slides', '_rels', src_slide_name + '.rels')
#     dst_slide_rels = os.path.join(temp_dst, 'ppt', 'slides', '_rels', new_slide_name + '.rels')
#     if os.path.exists(src_slide_rels):
#         shutil.copy(src_slide_rels, dst_slide_rels)
#
#     # Find layout from slide rels
#     rels_tree = ET.parse(src_slide_rels)
#     rels_root = rels_tree.getroot()
#     layout_target = None
#     for rel in rels_root.findall('Relationship'):
#         if rel.attrib['Type'].endswith('/slideLayout'):
#             layout_target = rel.attrib['Target'].replace('..', '').lstrip('/')
#             break
#
#     if layout_target:
#         layout_name = os.path.basename(layout_target)
#         layout_src = os.path.join(temp_src, 'ppt', 'slideLayouts', layout_name)
#         layout_dst = os.path.join(temp_dst, 'ppt', 'slideLayouts', layout_name)
#         shutil.copy(layout_src, layout_dst)
#
#         # Copy layout rels
#         layout_rels_src = os.path.join(temp_src, 'ppt', 'slideLayouts', '_rels', layout_name + '.rels')
#         layout_rels_dst = os.path.join(temp_dst, 'ppt', 'slideLayouts', '_rels', layout_name + '.rels')
#         if os.path.exists(layout_rels_src):
#             shutil.copy(layout_rels_src, layout_rels_dst)
#
#     # Modify presentation.xml
#     pres_path = os.path.join(temp_dst, 'ppt', 'presentation.xml')
#     pres_tree = ET.parse(pres_path)
#     pres_root = pres_tree.getroot()
#     sldIdLst = pres_root.find('{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
#     last_id = max([int(sld.get('id')) for sld in sldIdLst.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')])
#     new_id = last_id + 1
#     new_rid = f"rId{new_id + 100}"
#
#     new_elem = ET.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId', {
#         'id': str(new_id),
#         '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': new_rid
#     })
#     sldIdLst.append(new_elem)
#     pres_tree.write(pres_path, xml_declaration=True, encoding='utf-8')
#
#     # Add relationship
#     rels_path = os.path.join(temp_dst, 'ppt', '_rels', 'presentation.xml.rels')
#     rels_tree = ET.parse(rels_path)
#     rels_root = rels_tree.getroot()
#     new_rel = ET.Element('Relationship', {
#         'Id': new_rid,
#         'Type': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide',
#         'Target': f'slides/{new_slide_name}'
#     })
#     rels_root.append(new_rel)
#     rels_tree.write(rels_path, xml_declaration=True, encoding='utf-8')
#
#     # Rezip destination
#     new_pptx = dst_path.replace('.pptx', '_with_slide.pptx')
#     with zipfile.ZipFile(new_pptx, 'w', zipfile.ZIP_DEFLATED) as pptx:
#         for folder_name, subfolders, filenames in os.walk(temp_dst):
#             for filename in filenames:
#                 abs_path = os.path.join(folder_name, filename)
#                 rel_path = os.path.relpath(abs_path, temp_dst)
#                 pptx.write(abs_path, rel_path)
#
#     # Cleanup
#     shutil.rmtree(temp_src)
#     shutil.rmtree(temp_dst)
#
#     print(f"✅ اسلاید دوم با layout به فایل '{new_pptx}' منتقل شد.")

# Example usage:
# copy_slide_with_layout("template_3.pptx", "output_3.pptx", slide_index=2)


# import zipfile
# import os
# import shutil
# import tempfile
# import xml.etree.ElementTree as ET
#
#
# def create_pptx_with_single_slide(src_path, slide_index=1):
#     ET.register_namespace('', "http://schemas.openxmlformats.org/presentationml/2006/main")
#     ET.register_namespace('r', "http://schemas.openxmlformats.org/officeDocument/2006/relationships")
#
#     temp_src = tempfile.mkdtemp()
#     temp_dst = tempfile.mkdtemp()
#
#     # Unzip source pptx
#     with zipfile.ZipFile(src_path, 'r') as zip_ref:
#         zip_ref.extractall(temp_src)
#
#     # Create minimal pptx structure in destination
#     os.makedirs(os.path.join(temp_dst, '_rels'))
#     os.makedirs(os.path.join(temp_dst, 'docProps'))
#     os.makedirs(os.path.join(temp_dst, 'ppt', 'slides', '_rels'))
#     os.makedirs(os.path.join(temp_dst, 'ppt', '_rels'))
#     os.makedirs(os.path.join(temp_dst, 'ppt', 'slideLayouts', '_rels'))
#     os.makedirs(os.path.join(temp_dst, 'ppt', 'slideMasters', '_rels'))
#     os.makedirs(os.path.join(temp_dst, 'ppt', 'theme'))
#
#     # Copy core ppt files (docProps, content types)
#     for name in ['[Content_Types].xml']:
#         shutil.copy(os.path.join(temp_src, name), os.path.join(temp_dst, name))
#
#     for folder in ['docProps', 'ppt/theme']:
#         src_folder = os.path.join(temp_src, folder)
#         dst_folder = os.path.join(temp_dst, folder)
#         for file in os.listdir(src_folder):
#             src_file = os.path.join(src_folder, file)
#             if os.path.isfile(src_file):
#                 shutil.copy(src_file, os.path.join(dst_folder, file))
#
#     # Copy slide and rels
#     slide_id = slide_index + 1
#     slide_name = f'slide1.xml'
#     src_slide_name = f'slide{slide_id}.xml'
#
#     shutil.copy(os.path.join(temp_src, 'ppt', 'slides', src_slide_name),
#                 os.path.join(temp_dst, 'ppt', 'slides', slide_name))
#
#     # Copy slide rels
#     src_slide_rels = os.path.join(temp_src, 'ppt', 'slides', '_rels', f'{src_slide_name}.rels')
#     if os.path.exists(src_slide_rels):
#         shutil.copy(src_slide_rels,
#                     os.path.join(temp_dst, 'ppt', 'slides', '_rels', f'{slide_name}.rels'))
#
#     # Find layout
#     rels_tree = ET.parse(src_slide_rels)
#     rels_root = rels_tree.getroot()
#     layout_target = None
#     for rel in rels_root.findall('Relationship'):
#         if rel.attrib['Type'].endswith('/slideLayout'):
#             layout_target = rel.attrib['Target'].replace('..', '').lstrip('/')
#             break
#
#     # Copy slide layout and rels
#     if layout_target:
#         layout_name = os.path.basename(layout_target)
#         shutil.copy(os.path.join(temp_src, 'ppt', 'slideLayouts', layout_name),
#                     os.path.join(temp_dst, 'ppt', 'slideLayouts', layout_name))
#
#         rels_src = os.path.join(temp_src, 'ppt', 'slideLayouts', '_rels', layout_name + '.rels')
#         rels_dst = os.path.join(temp_dst, 'ppt', 'slideLayouts', '_rels', layout_name + '.rels')
#         if os.path.exists(rels_src) and os.path.isfile(rels_src):
#             shutil.copy(rels_src, rels_dst)
#
#     # Copy slide master and rels
#     for file in os.listdir(os.path.join(temp_src, 'ppt', 'slideMasters')):
#         src_file = os.path.join(temp_src, 'ppt', 'slideMasters', file)
#         if os.path.isfile(src_file):
#             shutil.copy(src_file, os.path.join(temp_dst, 'ppt', 'slideMasters', file))
#
#     src_rels_folder = os.path.join(temp_src, 'ppt', 'slideMasters', '_rels')
#     dst_rels_folder = os.path.join(temp_dst, 'ppt', 'slideMasters', '_rels')
#     for file in os.listdir(src_rels_folder):
#         src_file = os.path.join(src_rels_folder, file)
#         if os.path.isfile(src_file):
#             shutil.copy(src_file, os.path.join(dst_rels_folder, file))
#
#     # Copy presentation.xml and rels
#     shutil.copy(os.path.join(temp_src, 'ppt', 'presentation.xml'),
#                 os.path.join(temp_dst, 'ppt', 'presentation.xml'))
#     shutil.copy(os.path.join(temp_src, 'ppt', '_rels', 'presentation.xml.rels'),
#                 os.path.join(temp_dst, 'ppt', '_rels', 'presentation.xml.rels'))
#
#     # Copy root _rels
#     shutil.copy(os.path.join(temp_src, '_rels', '.rels'), os.path.join(temp_dst, '_rels', '.rels'))
#
#     # Copy media if exists
#     media_src = os.path.join(temp_src, 'ppt', 'media')
#     if os.path.exists(media_src):
#         media_dst = os.path.join(temp_dst, 'ppt', 'media')
#         shutil.copytree(media_src, media_dst)
#
#     # Zip to new pptx
#     output_pptx = src_path.replace('.pptx', '_single_slide.pptx')
#     with zipfile.ZipFile(output_pptx, 'w', zipfile.ZIP_DEFLATED) as pptx:
#         for folder_name, subfolders, filenames in os.walk(temp_dst):
#             for filename in filenames:
#                 abs_path = os.path.join(folder_name, filename)
#                 rel_path = os.path.relpath(abs_path, temp_dst)
#                 pptx.write(abs_path, rel_path)
#
#     shutil.rmtree(temp_src)
#     shutil.rmtree(temp_dst)
#
#     print(f"✅ فایل جدید با یک اسلاید ساخته شد: {output_pptx}")
#
#
# # Example usage
# create_pptx_with_single_slide("template_3.pptx", slide_index=2)  # اسلاید دوم




