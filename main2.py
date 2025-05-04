from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.autoshape import TextFrame, AutoShapeType
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


from lxml import etree

# فایل template.pptx را بارگذاری می‌کنیم
# template_path = 'template.pptx'
# prs = Presentation(template_path)
#
# # برای هر اسلاید در فایل template
# for slide in prs.slides:
#     for shape in slide.shapes:
#         if isinstance(shape, TextFrame):
#             print("این یک TextBox است.")
#         elif isinstance(shape, Picture):
#             print("این یک Picture است.")
#         elif isinstance(shape, SlidePlaceholder):
#             print("این یک SlidePlaceholder است.")
#         elif isinstance(shape, AutoShapeType):
#             print("این یک AutoShape است.")
#         else:
#             print(f"نوع ناشناخته: {shape.__class__}")
#
# # فایل جدید با نام جدید ذخیره می‌کنیم
# output_path = 'output.pptx'
# prs.save(output_path)

questionnaire = {'title':'مقایسه دو پلتفرم اسنپ و تپسی',
                 'questions':[{'text': 'کدام بهتر است؟'},{'text': 'اسنپ را چقدر دوست دارید؟'}]}
def clone_element(el):
    return etree.fromstring(etree.tostring(el))

class PptxElementDetector:
    elements = ['Picture','TextFrame','AutoShapeType','SlidePlaceholder', 'Shape']
    template_text = {
        'first_page': 'گزارش نظرسنجی {} در سنجاپ',
        'question_header': 'نظرسنجی انجام شده {} در سنجاپ'
    }
    def __init__(self, template_path):
        self.template_path = template_path
        self.prs = Presentation(template_path)
        # self.new_prs = Presentation()

    def get_first_slide(self):
        return self.prs.slides[0]

    def get_questions_slide(self):
        return self.prs.slides[1]


    def first_slide_layout(self):
        return self.prs.slides[0].slide_layout

    def last_slide_layout(self):
        return self.prs.slides[-1].slide_layout

    def question_slide_layout(self):
        return self.prs.slides[2].slide_layout

    def add_first_slide(self,title):
        new_slide = self.prs.slides.add_slide(self.first_slide_layout())
        for shape in self.prs.slides[0].shapes:
            # el = shape.element
            # new_el = el.clone()
            # new_slide.shapes._spTree.insert_element_before(clone_element(shape.element), 'p:extLst')
            if shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                self.set_text(new_shape, title)
                new_shape.text_frame.text = title

        self.add_question_slides()
        self.add_last_slide()
        output_path = 'output/output.pptx'
        self.prs.save(output_path)

    def add_question_slides(self):
        for question in questionnaire['questions']:
            new_slide = self.prs.slides.add_slide(self.question_slide_layout())
            for shape in self.get_questions_slide().shapes:

                if shape.has_text_frame:
                    new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)

                    if "?" in shape.text or "؟" in shape.text:
                        new_shape.text_frame.text = question['text']
                    else:
                        new_shape.text_frame.text = shape.text

    def shape_details(self):
        for shape in self.prs.slides[1].shapes:
            if (shape.has_chart):
                print(shape)


    def add_last_slide(self):
        new_slide = self.prs.slides.add_slide(self.last_slide_layout())
        for shape in self.prs.slides[2].shapes:
            # el = shape.element
            # new_el = el.clone()
            # new_slide.shapes._spTree.insert_element_before(clone_element(shape.element), 'p:extLst')
            if shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                # self.set_text(new_shape, title)
                # new_shape.text_frame.text = title

    def set_text(self, shape, text):
        if hasattr(shape, 'text'):
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    font_name = font.name
                    font_size = font.size
                    font_bold = font.bold
                    font_italic = font.italic

                    run.text = text

                    run.font.name = font_name
                    run.font.size = font_size
                    run.font.bold = font_bold
                    run.font.italic = font_italic


# result = PptxElementDetector('template_2.pptx').add_first_slide(questionnaire['title'])
# result = PptxElementDetector('template_2.pptx').shape_details()
result = PptxElementDetector('./v3/template_4.pptx').shape_details()
