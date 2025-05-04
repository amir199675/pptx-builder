from pptx import Presentation
from pptx.util import Inches, Pt

# ساختار نظرسنجی
questionnaire = {'title':'مقایسه دو پلتفرم اسنپ و تپسی',
                 'questions':[{'text': 'کدام بهتر است؟'},{'text': 'اسنپ را چقدر دوست دارید؟'}]}

# باز کردن تمپلیت
prs = Presentation('template.pptx')

# چاپ تعداد طرح‌بندی‌های موجود
print(f"تعداد طرح‌بندی‌های موجود: {len(prs.slide_layouts)}")

# به‌روزرسانی صفحه اول (عنوان)
title_slide = prs.slides[0]
for shape in title_slide.shapes:
    if shape.has_text_frame:
        shape.text = questionnaire['title']
        break

# به‌روزرسانی صفحه دوم برای هر سوال
for i, question in enumerate(questionnaire['questions']):
    if i == 0:  # استفاده از صفحه دوم تمپلیت
        question_slide = prs.slides[1]
    else:  # ایجاد صفحات جدید برای سوالات بعدی
        # استفاده از طرح‌بندی صفحه دوم برای اسلایدهای جدید
        question_slide = prs.slides.add_slide(prs.slides[1].slide_layout)
    
    # به‌روزرسانی متن سوال
    for shape in question_slide.shapes:
        if shape.has_text_frame:
            shape.text = f"سوال {i+1}: {question['text']}"
            break

# به‌روزرسانی صفحه آخر
last_slide = prs.slides[-1]
for shape in last_slide.shapes:
    if shape.has_text_frame:
        shape.text = "پایان نظرسنجی"
        break

# ذخیره پاورپوینت جدید
prs.save('output.pptx')